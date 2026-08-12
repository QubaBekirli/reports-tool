import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
import os
from app.config import settings


def extract_text_from_pdf(filepath: str) -> str:
    text_parts = []
    doc = fitz.open(filepath)
    for page in doc:
        page_text = page.get_text()
        if page_text.strip():
            text_parts.append(page_text)
        else:
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            try:
                ocr_text = pytesseract.image_to_string(img, lang="aze+eng")
            except Exception:
                try:
                    ocr_text = pytesseract.image_to_string(img, lang="eng")
                except Exception:
                    ocr_text = ""
            text_parts.append(ocr_text)
    doc.close()
    return "\n\n".join(text_parts)


def extract_text_from_docx(filepath: str) -> str:
    doc = DocxDocument(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n\n".join(parts)


def extract_text_from_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n\n".join(parts)


def extract_text_from_xlsx(filepath: str) -> str:
    wb = load_workbook(filepath, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n\n".join(parts)


def extract_text_from_image(filepath: str) -> str:
    img = Image.open(filepath)
    try:
        return pytesseract.image_to_string(img, lang="aze+eng")
    except Exception:
        return pytesseract.image_to_string(img, lang="eng")


def extract_text_from_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


EXTRACTORS = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "xlsx": extract_text_from_xlsx,
    "png": extract_text_from_image,
    "jpg": extract_text_from_image,
    "jpeg": extract_text_from_image,
    "txt": extract_text_from_txt,
}


def extract_text(filepath: str, format: str) -> str:
    extractor = EXTRACTORS.get(format)
    if not extractor:
        raise ValueError(f"Dəstəklənməyən format: {format}")
    return extractor(filepath)


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> list[str]:
    chunk_size = chunk_size or settings.chunk_size
    overlap = overlap or settings.chunk_overlap
    if overlap >= chunk_size:
        overlap = chunk_size // 2
    words = text.split()
    if len(words) <= chunk_size:
        return [text] if text.strip() else []
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end >= len(words):
            break
        start = end - overlap
    return chunks
